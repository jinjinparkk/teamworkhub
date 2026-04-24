"""
TeamWorkHub Premium Presentation Generator v4
White + Red elegant theme, Apple Keynote style — Korean version
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ─── Constants ────────────────────────────────────────────────────────
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
LIGHT_ROSE = RGBColor(0xFF, 0xF1, 0xF2)
ROSE_100 = RGBColor(0xFF, 0xE4, 0xE6)
ROSE_200 = RGBColor(0xFE, 0xCD, 0xD3)
RED_600 = RGBColor(0xDC, 0x26, 0x26)
RED_700 = RGBColor(0xB9, 0x1C, 0x1C)
RED_800 = RGBColor(0x99, 0x1B, 0x1B)
RED_900 = RGBColor(0x7F, 0x1D, 0x1D)
DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)
DARK_RED_BG = RGBColor(0x45, 0x0A, 0x0A)
CHARCOAL = RGBColor(0x11, 0x18, 0x27)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MID_TEXT = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
BLACK = RGBColor(0x00, 0x00, 0x00)
TRANSPARENT = None

# Code colors
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_GREEN = RGBColor(0xA6, 0xE3, 0xA1)
CODE_BLUE = RGBColor(0x89, 0xB4, 0xFA)
CODE_YELLOW = RGBColor(0xF9, 0xE2, 0xAF)
CODE_MAUVE = RGBColor(0xCB, 0xA6, 0xF7)
CODE_RED = RGBColor(0xF3, 0x8B, 0xA8)
CODE_TEXT = RGBColor(0xCD, 0xD6, 0xF4)
CODE_COMMENT = RGBColor(0x6C, 0x70, 0x86)

# Fonts
FONT_KR = "맑은 고딕"
FONT_CODE = "Consolas"

# Layout
MARGIN_LEFT = Inches(0.9)
MARGIN_TOP = Inches(0.7)
CONTENT_W = Inches(11.5)

# ─── Helper Functions ─────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_KR,
                font_size=Pt(14), font_color=MID_TEXT, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """Add a text box with one run of text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    # Set line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is not None:
        pPr.remove(lnSpc)
    lnSpc_el = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc_el.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc_el.append(spcPct)
    pPr.append(lnSpc_el)
    # Anchor
    txBox.text_frame.paragraphs[0].alignment = alignment
    return txBox


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a shape with optional fill and line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=Pt(0.5)):
    """Add a rectangle."""
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height,
                     fill_color, line_color, line_width)


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     line_color=None, line_width=Pt(0.5)):
    """Add a rounded rectangle."""
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                     fill_color, line_color, line_width)


def add_circle(slide, left, top, size, fill_color=None, line_color=None):
    """Add a circle (oval)."""
    return add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size,
                     fill_color, line_color)


def add_line(slide, left, top, width, color=RED_600, line_width=Pt(2)):
    """Add a horizontal line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.height = line_width
    return shape


def add_section_label(slide, left, top, text, color=RED_600):
    """Add a small section label."""
    tb = add_textbox(slide, left, top, Inches(4), Inches(0.4), text,
                     font_size=Pt(12), font_color=color, bold=True)
    # Add letter spacing via XML
    run = tb.text_frame.paragraphs[0].runs[0]
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', '300')
    return tb


def add_title_with_line(slide, left, top, title_text, font_size=Pt(32),
                        line_width=Inches(2.2)):
    """Add a title with a red accent line underneath."""
    tb = add_textbox(slide, left, top, Inches(10), Inches(0.6), title_text,
                     font_size=font_size, font_color=DARK_TEXT, bold=True,
                     line_spacing=1.1)
    add_line(slide, left, top + Inches(0.65), line_width)
    return tb


def add_card(slide, left, top, width, height, fill=WHITE, border=GRAY_200,
             accent_top=True, accent_color=RED_600, border_width=Pt(0.75)):
    """Add a card shape with optional red top accent."""
    card = add_rounded_rect(slide, left, top, width, height,
                            fill_color=fill, line_color=border,
                            line_width=border_width)
    if accent_top:
        add_rect(slide, left + Inches(0.05), top, width - Inches(0.1), Pt(3),
                 fill_color=accent_color)
    return card


def add_bullet_item(slide, left, top, text, font_size=Pt(14), text_color=MID_TEXT,
                    bullet_color=RED_600, spacing=Inches(0.38)):
    """Add a bullet point with a red circle bullet."""
    add_circle(slide, left, top + Inches(0.06), Inches(0.08), fill_color=bullet_color)
    add_textbox(slide, left + Inches(0.22), top, Inches(9), Inches(0.35), text,
                font_size=font_size, font_color=text_color, line_spacing=1.3)
    return top + spacing


def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.4),
                    color=RED_600):
    """Add a right-pointing chevron arrow."""
    shape = add_shape(slide, MSO_SHAPE.CHEVRON, left, top, width, height,
                      fill_color=color)
    return shape


def add_flow_box(slide, left, top, width, height, text, subtitle="",
                 fill=WHITE, border=GRAY_200, text_color=DARK_TEXT,
                 icon_text="", icon_color=RED_600):
    """Add a flow diagram box with icon and text."""
    card = add_rounded_rect(slide, left, top, width, height,
                            fill_color=fill, line_color=border, line_width=Pt(1))
    # Icon circle
    if icon_text:
        cx = left + width / 2 - Inches(0.25)
        cy = top + Inches(0.2)
        add_circle(slide, cx, cy, Inches(0.5), fill_color=icon_color)
        add_textbox(slide, cx, cy + Inches(0.05), Inches(0.5), Inches(0.4),
                    icon_text, font_size=Pt(16), font_color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
    # Main text
    ty = top + Inches(0.8) if icon_text else top + Inches(0.15)
    add_textbox(slide, left, ty, width, Inches(0.35), text,
                font_size=Pt(14), font_color=text_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Subtitle
    if subtitle:
        add_textbox(slide, left + Inches(0.1), ty + Inches(0.32), width - Inches(0.2),
                    Inches(0.5), subtitle,
                    font_size=Pt(10), font_color=GRAY_500,
                    alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    return card


def add_metric_card(slide, left, top, width, height, number, unit, description,
                    num_color=RED_600):
    """Add a metric display card."""
    card = add_card(slide, left, top, width, height, accent_color=num_color)
    add_textbox(slide, left, top + Inches(0.25), width, Inches(0.6), number,
                font_size=Pt(42), font_color=num_color, bold=True,
                alignment=PP_ALIGN.CENTER, line_spacing=1.0)
    add_textbox(slide, left, top + Inches(0.8), width, Inches(0.3), unit,
                font_size=Pt(14), font_color=GRAY_500, bold=False,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(1.1), width, Inches(0.5), description,
                font_size=Pt(11), font_color=GRAY_400,
                alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    return card


def set_shape_shadow(shape, blur=Pt(6), offset=Pt(2), color=BLACK, alpha=15):
    """Add a subtle drop shadow to a shape (via XML)."""
    spPr = shape._element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        spPr.append(effectLst)
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur)),
        'dist': str(int(offset)),
        'dir': '5400000',
        'algn': 'tl',
        'rotWithShape': '0'
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {
        'val': str(color) if isinstance(color, RGBColor) else '000000'
    })
    alphaEl = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
    srgbClr.append(alphaEl)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)


# ─── Presentation Setup ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Subtle gradient overlay — large dark red circle (decorative)
deco_circle = add_circle(slide, Inches(7.5), Inches(-2), Inches(9),
                         fill_color=RGBColor(0x2D, 0x0A, 0x0A))
# Make it semi-transparent via XML
sp = deco_circle._element.spPr
solidFill = sp.find(qn('a:solidFill'))
if solidFill is not None:
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '40000'})
        srgb.append(alpha)

# Small decorative circle top-left
add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(2.0),
           fill_color=RGBColor(0x35, 0x0A, 0x0A))

# Thin horizontal line accent
add_line(slide, Inches(0.9), Inches(2.3), Inches(1.5), RED_600, Pt(3))

# Title
add_textbox(slide, Inches(0.9), Inches(2.55), Inches(8), Inches(1.0),
            "옵대리",
            font_size=Pt(54), font_color=WHITE, bold=True, line_spacing=1.0)

# Subtitle
add_textbox(slide, Inches(0.9), Inches(3.5), Inches(8), Inches(0.8),
            "Gmail에서 Obsidian으로, 자동화.",
            font_size=Pt(22), font_color=GRAY_400, bold=False, line_spacing=1.2)

# Description line
add_textbox(slide, Inches(0.9), Inches(4.2), Inches(8), Inches(0.6),
            "AI 기반 이메일 인텔리전스 파이프라인",
            font_size=Pt(14), font_color=GRAY_500, line_spacing=1.3)

# Decorative dots (3 red dots)
for i in range(3):
    add_circle(slide, Inches(0.9) + Inches(i * 0.3), Inches(5.3),
               Inches(0.1), fill_color=RED_600)

# Footer
add_textbox(slide, Inches(0.9), Inches(6.3), Inches(10), Inches(0.35),
            "PM3팀, 박은진  |  2026.04",
            font_size=Pt(10), font_color=GRAY_500, line_spacing=1.0,
            bold=True)

# Version tag — top right
add_textbox(slide, Inches(10.5), Inches(0.5), Inches(2), Inches(0.3),
            "v1.0",
            font_size=Pt(11), font_color=GRAY_500, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.6), "PROBLEM")
add_title_with_line(slide, MARGIN_LEFT, Inches(1.0),
                    "이메일 관리, 왜 비효율적인가", Pt(30))

# 3 pain point cards
card_data = [
    {"icon": "01", "title": "시간 낭비",
     "desc": "매일 30분 이상을 이메일 분류,\n전달, 정리에 소비"},
    {"icon": "02", "title": "맥락 유실",
     "desc": "중요한 업무 지시가 긴 회신 체인\n(RE:RE:RE:)에 묻혀 유실"},
    {"icon": "03", "title": "가시성 부재",
     "desc": "팀 내 누가 무엇을 담당하는지\n한눈에 파악할 수 없음"},
]

card_w = Inches(3.3)
card_h = Inches(2.8)
gap = Inches(0.4)
total_cards_w = 3 * card_w + 2 * gap
start_x = (SLIDE_W - total_cards_w) // 2
card_y = Inches(2.2)

for i, cd in enumerate(card_data):
    cx = start_x + i * (card_w + gap)
    card = add_card(slide, cx, card_y, card_w, card_h, fill=WHITE,
                    border=GRAY_200, accent_color=RED_600)
    set_shape_shadow(card, blur=Pt(8), offset=Pt(3), alpha=8)

    # Number circle
    num_circle = add_circle(slide, cx + Inches(0.2), card_y + Inches(0.3),
                            Inches(0.55), fill_color=LIGHT_ROSE)
    add_textbox(slide, cx + Inches(0.2), card_y + Inches(0.35),
                Inches(0.55), Inches(0.45), cd["icon"],
                font_size=Pt(18), font_color=RED_600, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Card title
    add_textbox(slide, cx + Inches(0.2), card_y + Inches(1.05),
                card_w - Inches(0.4), Inches(0.35), cd["title"],
                font_size=Pt(16), font_color=DARK_TEXT, bold=True)

    # Card description
    add_textbox(slide, cx + Inches(0.2), card_y + Inches(1.45),
                card_w - Inches(0.4), Inches(1.0), cd["desc"],
                font_size=Pt(11), font_color=GRAY_500, line_spacing=1.5)

# Bottom accent bar
add_rect(slide, Inches(0), Inches(7.0), Inches(13.33), Pt(4), fill_color=RED_600)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY AUTOMATE (Necessity + Benefits)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.5), "WHY AUTOMATE")
add_title_with_line(slide, MARGIN_LEFT, Inches(0.85),
                    "왜 자동화가 필요한가", Pt(30))

# ── Left column: 현재 상황 (dark) ──
left_x = Inches(0.9)
left_w = Inches(5.3)
col_top = Inches(1.8)

# Column header — dark card
left_header = add_rounded_rect(slide, left_x, col_top, left_w, Inches(0.5),
                                fill_color=DARK_BG, line_color=DARK_BG, line_width=Pt(0))
add_textbox(slide, left_x + Inches(0.25), col_top + Inches(0.08),
            left_w - Inches(0.5), Inches(0.35), "현재 상황",
            font_size=Pt(14), font_color=WHITE, bold=True)

pain_items = [
    ("30건+", "하루 평균 이메일",
     "수동 분류·전달·정리에\n매일 30분 이상 소요"),
    ("RE:RE:", "복잡한 회신 체인",
     "다수 이해관계자 간 메일이\n수십 회 오가며 맥락 파편화"),
    ("0%", "업무 추적률",
     "담당자·마감일·액션아이템이\n개인 메일함에 분산되어 추적 불가"),
    ("수동", "보고서 작성",
     "업무 현황 파악을 위해\n매번 메일을 다시 뒤져야 함"),
]

for pi, (num, title, desc) in enumerate(pain_items):
    py = col_top + Inches(0.7) + pi * Inches(1.1)

    # Number badge
    badge = add_rounded_rect(slide, left_x + Inches(0.15), py + Inches(0.05),
                              Inches(0.85), Inches(0.35),
                              fill_color=GRAY_100, line_color=GRAY_200,
                              line_width=Pt(0.5))
    add_textbox(slide, left_x + Inches(0.15), py + Inches(0.06),
                Inches(0.85), Inches(0.3), num,
                font_size=Pt(13), font_color=GRAY_500, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, left_x + Inches(1.15), py,
                Inches(3.5), Inches(0.3), title,
                font_size=Pt(13), font_color=DARK_TEXT, bold=True)
    add_textbox(slide, left_x + Inches(1.15), py + Inches(0.32),
                Inches(3.8), Inches(0.6), desc,
                font_size=Pt(10), font_color=GRAY_500, line_spacing=1.4)

# ── Center divider arrow ──
div_x = Inches(6.35)
for ai in range(3):
    ay = Inches(2.8) + ai * Inches(1.4)
    add_arrow_right(slide, div_x, ay, Inches(0.4), Inches(0.35), RED_600)

# ── Right column: 자동화 도입 효과 (bright) ──
right_x = Inches(7.0)
right_w = Inches(5.3)

# Column header — red card
right_header = add_rounded_rect(slide, right_x, col_top, right_w, Inches(0.5),
                                 fill_color=RED_600, line_color=RED_600, line_width=Pt(0))
add_textbox(slide, right_x + Inches(0.25), col_top + Inches(0.08),
            right_w - Inches(0.5), Inches(0.35), "자동화 도입 시",
            font_size=Pt(14), font_color=WHITE, bold=True)

benefit_items = [
    ("자동", "출근 즉시 업무 파악",
     "매일 아침 자동 생성된 일일 노트로\n전날 메일 현황을 즉시 확인"),
    ("AI", "핵심만 추출, 맥락 유지",
     "회신 체인에서 최신 답장만 추출하고\nAI가 3줄 요약 + 담당자 자동 배정"),
    ("자산화", "업무 이력이 지식으로 축적",
     "모든 이메일이 마크다운으로 기록되어\n검색·링크·쿼리 가능한 팀 자산이 됨"),
    ("실시간", "팀 업무 가시성 확보",
     "담당자별 대시보드로\n업무 분배와 진행 현황 즉시 파악"),
]

for bi, (num, title, desc) in enumerate(benefit_items):
    by = col_top + Inches(0.7) + bi * Inches(1.1)

    # Number badge (red)
    badge = add_rounded_rect(slide, right_x + Inches(0.15), by + Inches(0.05),
                              Inches(0.85), Inches(0.35),
                              fill_color=LIGHT_ROSE, line_color=ROSE_200,
                              line_width=Pt(0.5))
    add_textbox(slide, right_x + Inches(0.15), by + Inches(0.06),
                Inches(0.85), Inches(0.3), num,
                font_size=Pt(13), font_color=RED_600, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, right_x + Inches(1.15), by,
                Inches(3.5), Inches(0.3), title,
                font_size=Pt(13), font_color=DARK_TEXT, bold=True)
    add_textbox(slide, right_x + Inches(1.15), by + Inches(0.32),
                Inches(3.8), Inches(0.6), desc,
                font_size=Pt(10), font_color=GRAY_500, line_spacing=1.4)

# Bottom accent bar
add_rect(slide, Inches(0), Inches(6.95), Inches(13.33), Pt(4), fill_color=RED_600)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION (Pipeline Flow)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.6), "SOLUTION")
add_title_with_line(slide, MARGIN_LEFT, Inches(1.0),
                    "자동화 파이프라인", Pt(30))

# Subtitle
add_textbox(slide, MARGIN_LEFT, Inches(1.85), Inches(10), Inches(0.4),
            "받은편지함에서 인사이트까지, 자동으로.",
            font_size=Pt(13), font_color=GRAY_500, bold=True)

# Flow diagram: 4 boxes with arrows
flow_items = [
    {"icon": "G", "title": "Gmail -> Drive", "sub": "이메일 수집"},
    {"icon": "AI", "title": "Claude AI", "sub": "분석 및 분류"},
    {"icon": "O", "title": "Obsidian", "sub": "노트 자동 생성"},
    {"icon": "R", "title": "리포트", "sub": "일간 리포트\n담당자별 체크리스트\n대시보드"},
]

box_w = Inches(2.2)
box_h = Inches(2.2)
arrow_w = Inches(0.45)
total_w = 4 * box_w + 3 * (arrow_w + Inches(0.2))
flow_start_x = (SLIDE_W - total_w) // 2
flow_y = Inches(2.8)

for i, item in enumerate(flow_items):
    bx = flow_start_x + i * (box_w + arrow_w + Inches(0.2))

    card = add_rounded_rect(slide, bx, flow_y, box_w, box_h,
                            fill_color=WHITE, line_color=GRAY_200,
                            line_width=Pt(1.2))
    set_shape_shadow(card, blur=Pt(10), offset=Pt(4), alpha=8)

    # Icon circle
    icon_size = Inches(0.7)
    icx = bx + box_w // 2 - icon_size // 2
    icy = flow_y + Inches(0.3)
    ic = add_circle(slide, icx, icy, icon_size, fill_color=RED_600)
    add_textbox(slide, icx, icy + Inches(0.1), icon_size, Inches(0.5),
                item["icon"], font_size=Pt(20), font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, bx, flow_y + Inches(1.15), box_w, Inches(0.35),
                item["title"], font_size=Pt(16), font_color=DARK_TEXT, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, bx + Inches(0.15), flow_y + Inches(1.55),
                box_w - Inches(0.3), Inches(0.55),
                item["sub"], font_size=Pt(11), font_color=GRAY_500,
                alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Arrow
    if i < 3:
        ax = bx + box_w + Inches(0.05)
        ay = flow_y + box_h // 2 - Inches(0.2)
        add_arrow_right(slide, ax, ay, Inches(0.35), Inches(0.4), RED_600)

# Bottom description bar
desc_bar_y = Inches(5.5)
add_rect(slide, Inches(0.9), desc_bar_y, Inches(11.5), Inches(0.9),
         fill_color=LIGHT_ROSE)
add_textbox(slide, Inches(1.3), desc_bar_y + Inches(0.15), Inches(10.7), Inches(0.6),
            "Cloud Scheduler로 매일 평일 09:00에 자동 실행. 수동 작업 제로.",
            font_size=Pt(12), font_color=RED_800, line_spacing=1.6)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — FEATURE: AI Analysis
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.6), "FEATURE 01")
add_title_with_line(slide, MARGIN_LEFT, Inches(1.0),
                    "AI 기반 이메일 인텔리전스", Pt(30))

# Left side: 4 capability cards (2x2)
cap_data = [
    {"title": "스마트 요약", "desc": "Claude AI가 이메일별\n한국어 요약을 자동 생성"},
    {"title": "담당자 자동 배정", "desc": "정규식 → Claude 추론 → 이메일\n3단계 폴백으로 담당자 감지"},
    {"title": "우선순위 분류", "desc": "긴급도 자동 분류:\n긴급 / 보통 / 낮음"},
    {"title": "카테고리 태깅", "desc": "자동 분류: 보고, 승인요청,\n공지, 미팅, 일반"},
]

mini_w = Inches(3.2)
mini_h = Inches(1.5)
mini_gap = Inches(0.25)
mini_start_x = Inches(0.9)
mini_start_y = Inches(2.0)

for i, cap in enumerate(cap_data):
    row = i // 2
    col = i % 2
    mx = mini_start_x + col * (mini_w + mini_gap)
    my = mini_start_y + row * (mini_h + mini_gap)

    card = add_card(slide, mx, my, mini_w, mini_h, fill=WHITE, border=GRAY_200)
    set_shape_shadow(card, blur=Pt(6), offset=Pt(2), alpha=6)

    # Red dot indicator
    add_circle(slide, mx + Inches(0.2), my + Inches(0.25), Inches(0.12),
               fill_color=RED_600)

    add_textbox(slide, mx + Inches(0.45), my + Inches(0.18), mini_w - Inches(0.6),
                Inches(0.3), cap["title"],
                font_size=Pt(14), font_color=DARK_TEXT, bold=True)

    add_textbox(slide, mx + Inches(0.45), my + Inches(0.55), mini_w - Inches(0.6),
                Inches(0.8), cap["desc"],
                font_size=Pt(11), font_color=GRAY_500, line_spacing=1.5)

# Right side: Reply chain handling box
rc_x = Inches(7.3)
rc_y = Inches(2.0)
rc_w = Inches(5.0)
rc_h = Inches(3.55)

rc_card = add_rounded_rect(slide, rc_x, rc_y, rc_w, rc_h,
                           fill_color=GRAY_100, line_color=GRAY_200,
                           line_width=Pt(1))

add_textbox(slide, rc_x + Inches(0.3), rc_y + Inches(0.2), rc_w - Inches(0.6),
            Inches(0.3), "회신 체인 추출",
            font_size=Pt(13), font_color=RED_600, bold=True)

# Simulated email chain
chain_lines = [
    ("RE:RE:RE: 대시보드 Task 리스트 공유", GRAY_400, False),
    ("RE:RE: 대시보드 Task 리스트 공유", GRAY_400, False),
    ("RE: 대시보드 Task 리스트 공유", GRAY_400, False),
    ("대시보드 Task 리스트 공유", RED_600, True),
]

for j, (line, color, is_bold) in enumerate(chain_lines):
    ly = rc_y + Inches(0.65) + j * Inches(0.4)
    prefix = "  " * j
    # Strike-through effect for old ones, highlight for latest
    if is_bold:
        add_rect(slide, rc_x + Inches(0.3), ly - Inches(0.05),
                 rc_w - Inches(0.6), Inches(0.35), fill_color=ROSE_100)
    add_textbox(slide, rc_x + Inches(0.4) + Inches(j * 0.15), ly,
                rc_w - Inches(1), Inches(0.3), line,
                font_size=Pt(11), font_color=color, bold=is_bold,
                font_name=FONT_CODE)

# Arrow pointing to extracted
add_textbox(slide, rc_x + Inches(0.3), rc_y + Inches(2.4), rc_w - Inches(0.6),
            Inches(0.8),
            "가장 최근의 의미 있는 내용만\n추출하여 AI가 처리합니다.",
            font_size=Pt(10), font_color=GRAY_500, line_spacing=1.5)

# Bottom info bar
add_rect(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.6),
         fill_color=DARK_BG)
add_textbox(slide, Inches(1.3), Inches(5.98), Inches(10.7), Inches(0.4),
            "Claude Haiku 4.5 기반  |  이메일당 API 1회 호출  |  처리 시간 2초 미만",
            font_size=Pt(11), font_color=GRAY_400, line_spacing=1.0)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE: Daily Note (Code Preview)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.5), "FEATURE 02")
add_title_with_line(slide, MARGIN_LEFT, Inches(0.85),
                    "자동 일일 노트", Pt(30))

add_textbox(slide, MARGIN_LEFT, Inches(1.65), Inches(5), Inches(0.35),
            "매일 아침, Obsidian에 완성된 일일 브리핑이 준비됩니다.",
            font_size=Pt(12), font_color=GRAY_500, italic=True)

# Code editor mockup
editor_x = Inches(0.9)
editor_y = Inches(2.2)
editor_w = Inches(6.8)
editor_h = Inches(4.5)

# Editor background
editor_bg = add_rounded_rect(slide, editor_x, editor_y, editor_w, editor_h,
                              fill_color=CODE_BG, line_color=RGBColor(0x31, 0x31, 0x45),
                              line_width=Pt(1))
set_shape_shadow(editor_bg, blur=Pt(15), offset=Pt(6), alpha=20)

# Title bar
add_rect(slide, editor_x, editor_y, editor_w, Inches(0.35),
         fill_color=RGBColor(0x18, 0x18, 0x28))

# Traffic light dots
dot_y = editor_y + Inches(0.12)
dot_colors = [RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFF, 0xBD, 0x2E),
              RGBColor(0x28, 0xCA, 0x41)]
for k, dc in enumerate(dot_colors):
    add_circle(slide, editor_x + Inches(0.2) + k * Inches(0.25), dot_y,
               Inches(0.1), fill_color=dc)

# File name in title bar
add_textbox(slide, editor_x + Inches(1.2), editor_y + Inches(0.05),
            Inches(3), Inches(0.28), "2026-04-21.md",
            font_name=FONT_CODE, font_size=Pt(10), font_color=GRAY_400,
            alignment=PP_ALIGN.CENTER)

# Code content — matches actual daily_writer output
code_lines = [
    ("---", CODE_COMMENT),
    ("Type: daily_note", CODE_BLUE),
    ("date: 2026-04-21", CODE_BLUE),
    ("period: 2026-04-20 18:00 ~ 04-21 09:00", CODE_BLUE),
    ("email_count: 5", CODE_BLUE),
    ("assignees: [박은진, 이해랑, 이자명]", CODE_GREEN),
    ("categories: [보고, 승인요청, 공지]", CODE_GREEN),
    ("has_urgent: true", CODE_RED),
    ("---", CODE_COMMENT),
    ("", None),
    ("### Today's work", CODE_YELLOW),
    ("#### To do list", CODE_MAUVE),
    ("- [ ] [[대시보드 Task 리스트 공유]] #이해랑", CODE_TEXT),
    ("  담당자:: 이해랑 | 우선순위:: 긴급 | 카테고리:: 보고", CODE_RED),
    ("- [ ] [[Affiliate Tracking 권한 요청]] #이자명", CODE_TEXT),
    ("  담당자:: 이자명 | 카테고리:: 승인요청", CODE_GREEN),
    ("- [x] [[데이터 검증 일일보고]] #박은진", CODE_COMMENT),
    ("", None),
    ("#### 정기적인 일", CODE_YELLOW),
    ("- [ ] RPA", CODE_TEXT),
    ("", None),
    ("### 미완료", CODE_MAUVE),
]

line_h = Inches(0.2)
code_y = editor_y + Inches(0.5)
for ci, (cline, ccolor) in enumerate(code_lines):
    if ccolor is None:
        continue
    add_textbox(slide, editor_x + Inches(0.3), code_y + ci * line_h,
                editor_w - Inches(0.6), Inches(0.22), cline,
                font_name=FONT_CODE, font_size=Pt(9), font_color=ccolor,
                line_spacing=1.0)

# Right side: Feature highlights
feat_x = Inches(8.1)
feat_y = Inches(2.2)
feat_items = [
    ("YAML 프론트매터", "Type, assignees, categories 등\nDataview 호환 메타데이터"),
    ("위키 링크", "[[제목]] 링크로\n개별 이메일 노트에 연결"),
    ("담당자 태그", "#이름 태그로\n담당자별 필터링 가능"),
    ("인라인 필드", "담당자:: / 우선순위:: / 카테고리::\nDataview 쿼리용"),
    ("정기적인 일", "월=RPA / 화=로직점검 등\n요일별 자동 생성"),
    ("야간 메일 수집", "18:00~08:59 메일 수집\n월요일은 금요일 18시부터"),
]

for fi, (ftitle, fdesc) in enumerate(feat_items):
    fy = feat_y + fi * Inches(0.72)
    add_circle(slide, feat_x, fy + Inches(0.02), Inches(0.1), fill_color=RED_600)
    add_textbox(slide, feat_x + Inches(0.22), fy - Inches(0.03), Inches(4), Inches(0.25),
                ftitle, font_size=Pt(12), font_color=DARK_TEXT, bold=True)
    add_textbox(slide, feat_x + Inches(0.22), fy + Inches(0.22), Inches(4), Inches(0.4),
                fdesc, font_size=Pt(9), font_color=GRAY_500, line_spacing=1.4)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURE: Reports (Weekly / Monthly / Dashboard)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.6), "FEATURE 03")
add_title_with_line(slide, MARGIN_LEFT, Inches(1.0),
                    "종합 리포팅 시스템", Pt(30))

add_textbox(slide, MARGIN_LEFT, Inches(1.85), Inches(10), Inches(0.35),
            "개별 이메일 노트, 일일 브리핑, 팀 대시보드까지 자동 생성.",
            font_size=Pt(12), font_color=GRAY_500, italic=True)

# Three report cards
report_data = [
    {
        "title": "이메일 노트",
        "subtitle": "{제목}.md",
        "items": [
            "YAML 메타데이터 (제목/발신자/태그)",
            "AI 요약 (3줄 핵심 요약)",
            "본문 자동 정리 (서명/인용 제거)",
            "첨부파일 Drive 링크 연결",
        ],
        "accent": RED_600,
    },
    {
        "title": "일일 노트",
        "subtitle": "YYYY-MM-DD.md",
        "items": [
            "위키링크 + 담당자 태그 체크리스트",
            "인라인 필드 (담당자/우선순위/카테고리)",
            "요일별 정기적인 일 자동 생성",
            "미완료 항목 Dataview 추적",
        ],
        "accent": RED_700,
    },
    {
        "title": "대시보드",
        "subtitle": "Dashboard.md + 담당자별.md",
        "items": [
            "긴급 메일 / 미처리 항목 쿼리",
            "담당자별 개별 페이지 자동 생성",
            "최근 7일 일간 다이제스트",
            "카테고리별 / 담당자별 통계",
        ],
        "accent": RED_800,
    },
]

rcard_w = Inches(3.5)
rcard_h = Inches(3.5)
rcard_gap = Inches(0.35)
rcard_start_x = Inches(0.9)
rcard_y = Inches(2.5)

for ri, rd in enumerate(report_data):
    rx = rcard_start_x + ri * (rcard_w + rcard_gap)
    card = add_card(slide, rx, rcard_y, rcard_w, rcard_h,
                    fill=WHITE, border=GRAY_200, accent_color=rd["accent"])
    set_shape_shadow(card, blur=Pt(8), offset=Pt(3), alpha=8)

    # Title
    add_textbox(slide, rx + Inches(0.3), rcard_y + Inches(0.25),
                rcard_w - Inches(0.6), Inches(0.35), rd["title"],
                font_size=Pt(18), font_color=DARK_TEXT, bold=True)

    # Subtitle (filename)
    add_textbox(slide, rx + Inches(0.3), rcard_y + Inches(0.6),
                rcard_w - Inches(0.6), Inches(0.25), rd["subtitle"],
                font_name=FONT_CODE, font_size=Pt(10), font_color=rd["accent"])

    # Divider
    add_line(slide, rx + Inches(0.3), rcard_y + Inches(0.95),
             rcard_w - Inches(0.6), GRAY_200, Pt(1))

    # Items
    for ii, item_text in enumerate(rd["items"]):
        iy = rcard_y + Inches(1.15) + ii * Inches(0.5)
        add_circle(slide, rx + Inches(0.35), iy + Inches(0.06),
                   Inches(0.07), fill_color=rd["accent"])
        add_textbox(slide, rx + Inches(0.55), iy, rcard_w - Inches(0.9),
                    Inches(0.4), item_text,
                    font_size=Pt(12), font_color=MID_TEXT, line_spacing=1.3)

# Bottom note
add_textbox(slide, MARGIN_LEFT, Inches(6.3), Inches(11), Inches(0.35),
            "모든 노트는 Dataview 인라인 필드가 포함된 Obsidian 호환 마크다운으로 생성됩니다.",
            font_size=Pt(10), font_color=GRAY_400, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE (System Diagram)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.5), "ARCHITECTURE")
add_title_with_line(slide, MARGIN_LEFT, Inches(0.85),
                    "시스템 구조", Pt(30))

# Architecture diagram — clean boxes and arrows
# Layer 1: Source
src_y = Inches(2.0)
src_h = Inches(1.0)

# Gmail box
gmail_x = Inches(0.9)
gmail_w = Inches(2.5)
add_flow_box(slide, gmail_x, src_y, gmail_w, src_h,
             "Gmail API", "OAuth 2.0 Desktop\nIME notifications",
             fill=WHITE, border=GRAY_200, icon_text="", icon_color=RED_600)
# Red left bar
add_rect(slide, gmail_x, src_y, Pt(4), src_h, fill_color=RED_600)

# Arrow down from Gmail
arr1_x = gmail_x + gmail_w // 2 - Inches(0.15)
arr1_y = src_y + src_h + Inches(0.1)
# Vertical connector
add_rect(slide, arr1_x + Inches(0.12), arr1_y, Pt(2.5), Inches(0.4), fill_color=RED_600)

# Layer 2: Processing (FastAPI)
proc_y = Inches(3.5)
proc_h = Inches(1.2)
proc_x = Inches(0.5)
proc_w = Inches(12.3)

proc_bg = add_rounded_rect(slide, proc_x, proc_y, proc_w, proc_h,
                           fill_color=LIGHT_ROSE, line_color=ROSE_200,
                           line_width=Pt(1))

add_textbox(slide, proc_x + Inches(0.3), proc_y + Inches(0.1),
            Inches(3), Inches(0.25), "FastAPI 애플리케이션 (Cloud Run)",
            font_size=Pt(11), font_color=RED_700, bold=True)

# Processing sub-boxes
proc_boxes = [
    ("gmail_client", "수집 &\n파싱"),
    ("summarizer", "Claude AI\n분석"),
    ("md_writer", "노트\n생성"),
    ("daily_writer", "일일\n노트"),
    ("dashboard", "대시보드\n+ 담당자"),
    ("drive_client", "Drive\n백업"),
]

pb_w = Inches(1.7)
pb_h = Inches(0.65)
pb_gap = Inches(0.15)
pb_start_x = proc_x + Inches(0.3)
pb_y = proc_y + Inches(0.4)

for pi, (pname, pdesc) in enumerate(proc_boxes):
    px = pb_start_x + pi * (pb_w + pb_gap)
    pb = add_rounded_rect(slide, px, pb_y, pb_w, pb_h,
                          fill_color=WHITE, line_color=ROSE_200,
                          line_width=Pt(0.75))
    add_textbox(slide, px + Inches(0.08), pb_y + Inches(0.05),
                pb_w - Inches(0.16), Inches(0.2), pname,
                font_name=FONT_CODE, font_size=Pt(8), font_color=RED_700,
                bold=True)
    add_textbox(slide, px + Inches(0.08), pb_y + Inches(0.28),
                pb_w - Inches(0.16), Inches(0.35), pdesc,
                font_size=Pt(8), font_color=GRAY_500, line_spacing=1.2)

    # Arrows between boxes
    if pi < len(proc_boxes) - 1:
        ax = px + pb_w + Inches(0.02)
        ay = pb_y + pb_h // 2 - Inches(0.1)
        add_arrow_right(slide, ax, ay, Inches(0.12), Inches(0.2), RED_600)

# Layer 3: Output
out_y = Inches(5.1)
out_h = Inches(0.85)

# Vertical connectors from processing to outputs
for oi in range(3):
    cx = Inches(3.0) + oi * Inches(3.5)
    add_rect(slide, cx, proc_y + proc_h, Pt(2.5), Inches(0.35), fill_color=RED_600)

# Output boxes
out_data = [
    ("Google Drive", "클라우드 백업\ntwh_{msgId}.md"),
    ("Obsidian Vault", "로컬 마크다운\n연결된 노트"),
    ("대시보드", "Dataview 쿼리\n담당자 페이지"),
]

out_w = Inches(2.8)
out_gap = Inches(0.7)
out_start_x = Inches(1.5)

for oi, (otitle, odesc) in enumerate(out_data):
    ox = out_start_x + oi * (out_w + out_gap)
    ob = add_rounded_rect(slide, ox, out_y, out_w, out_h,
                          fill_color=WHITE, line_color=GRAY_200,
                          line_width=Pt(1))
    add_rect(slide, ox, out_y + out_h - Pt(4), out_w, Pt(4), fill_color=RED_600)
    add_textbox(slide, ox + Inches(0.15), out_y + Inches(0.1),
                out_w - Inches(0.3), Inches(0.25), otitle,
                font_size=Pt(12), font_color=DARK_TEXT, bold=True)
    add_textbox(slide, ox + Inches(0.15), out_y + Inches(0.38),
                out_w - Inches(0.3), Inches(0.4), odesc,
                font_name=FONT_CODE, font_size=Pt(9), font_color=GRAY_500,
                line_spacing=1.3)

# Side label: Cloud Scheduler
add_rect(slide, Inches(11.5), Inches(2.0), Inches(1.6), Inches(1.0),
         fill_color=DARK_BG)
add_textbox(slide, Inches(11.5), Inches(2.15), Inches(1.6), Inches(0.25),
            "Cloud", font_size=Pt(10), font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(11.5), Inches(2.4), Inches(1.6), Inches(0.25),
            "Scheduler", font_size=Pt(10), font_color=GRAY_400,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(11.5), Inches(2.7), Inches(1.6), Inches(0.25),
            "09:00 KST", font_size=Pt(9), font_color=RED_600, bold=True,
            alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — IMPACT (Before/After + Metrics)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, OFF_WHITE)

add_section_label(slide, MARGIN_LEFT, Inches(0.6), "IMPACT")
add_title_with_line(slide, MARGIN_LEFT, Inches(1.0),
                    "측정 가능한 개선 효과", Pt(30))

# Before / After comparison
ba_y = Inches(2.0)
ba_h = Inches(2.0)

# Before card
before_x = Inches(0.9)
before_w = Inches(5.4)
before_card = add_card(slide, before_x, ba_y, before_w, ba_h,
                       fill=WHITE, border=GRAY_200, accent_color=GRAY_400)
set_shape_shadow(before_card, blur=Pt(6), offset=Pt(2), alpha=6)

add_textbox(slide, before_x + Inches(0.3), ba_y + Inches(0.2),
            Inches(2), Inches(0.3), "도입 전",
            font_size=Pt(14), font_color=GRAY_400, bold=True)

before_items = [
    "매일 30건 이상 이메일 수동 분류",
    "회신 체인에서 업무 지시 유실",
    "수작업 주간 보고서 작성 (2시간 이상)",
    "팀 업무 분배 현황 파악 불가",
]
for bi, btext in enumerate(before_items):
    by = ba_y + Inches(0.65) + bi * Inches(0.32)
    add_textbox(slide, before_x + Inches(0.35), by,
                Inches(0.2), Inches(0.25), "x",
                font_name=FONT_CODE, font_size=Pt(12), font_color=GRAY_400, bold=True)
    add_textbox(slide, before_x + Inches(0.6), by,
                before_w - Inches(0.9), Inches(0.25), btext,
                font_size=Pt(11), font_color=GRAY_500, line_spacing=1.2)

# Arrow between
add_arrow_right(slide, Inches(6.5), ba_y + ba_h // 2 - Inches(0.2),
                Inches(0.5), Inches(0.4), RED_600)

# After card
after_x = Inches(7.2)
after_w = Inches(5.4)
after_card = add_card(slide, after_x, ba_y, after_w, ba_h,
                      fill=WHITE, border=GRAY_200, accent_color=RED_600)
set_shape_shadow(after_card, blur=Pt(6), offset=Pt(2), alpha=6)

add_textbox(slide, after_x + Inches(0.3), ba_y + Inches(0.2),
            Inches(2), Inches(0.3), "도입 후",
            font_size=Pt(14), font_color=RED_600, bold=True)

after_items = [
    "전체 이메일 자동 처리 및 분류",
    "스마트 회신 체인 추출, 누락 제로",
    "리포트 자동 생성 (노력 제로)",
    "담당자 추적 팀 대시보드 완비",
]
for ai_idx, atext in enumerate(after_items):
    ay = ba_y + Inches(0.65) + ai_idx * Inches(0.32)
    add_textbox(slide, after_x + Inches(0.3), ay,
                Inches(0.25), Inches(0.25), "+",
                font_name=FONT_CODE, font_size=Pt(14), font_color=RED_600, bold=True)
    add_textbox(slide, after_x + Inches(0.6), ay,
                after_w - Inches(0.9), Inches(0.25), atext,
                font_size=Pt(11), font_color=MID_TEXT, line_spacing=1.2)

# 4 Metric cards at bottom
metrics = [
    {"num": "30", "unit": "분 절약/일", "desc": "이메일 처리시간 제거"},
    {"num": "0", "unit": "누락 항목", "desc": "모든 업무 지시\n완전 수집"},
    {"num": "100%", "unit": "기록률", "desc": "모든 이메일 자동 저장\n및 검색 가능"},
    {"num": "1-click", "unit": "리포트", "desc": "일간, 주간, 월간 즉시 생성"},
]

met_w = Inches(2.6)
met_h = Inches(1.9)
met_gap = Inches(0.3)
met_start_x = Inches(0.9)
met_y = Inches(4.5)

for mi, m in enumerate(metrics):
    mx = met_start_x + mi * (met_w + met_gap)
    mcard = add_metric_card(slide, mx, met_y, met_w, met_h,
                            m["num"], m["unit"], m["desc"])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — NEXT STEPS + THANK YOU
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Decorative circle
add_circle(slide, Inches(-1), Inches(3), Inches(5),
           fill_color=RGBColor(0x2D, 0x0A, 0x0A))

# Small accent circles
add_circle(slide, Inches(11), Inches(0.5), Inches(0.8),
           fill_color=RGBColor(0x35, 0x0A, 0x0A))
add_circle(slide, Inches(11.8), Inches(1.0), Inches(0.4),
           fill_color=RGBColor(0x45, 0x0A, 0x0A))

# Section label
add_section_label(slide, Inches(0.9), Inches(0.7), "NEXT STEPS", RED_600)

# Title
add_textbox(slide, Inches(0.9), Inches(1.1), Inches(6), Inches(0.6),
            "로드맵 및 향후 계획",
            font_size=Pt(32), font_color=WHITE, bold=True)

add_line(slide, Inches(0.9), Inches(1.8), Inches(2), RED_600, Pt(3))

# Roadmap items
roadmap = [
    ("완료", "Cloud Run 배포", "GCP 배포 + Cloud Scheduler로\n매일 평일 08:00 자동 실행 중"),
    ("진행중", "업무 자산화", "모든 이메일·회의·업무 기록을\nObsidian 마크다운으로 축적"),
    ("예정", "Claude 연동", "축적된 업무 자산을 Claude AI와 연결하여\n과거 이력 검색·분석·질의응답"),
    ("확장", "팀 지식 허브", "이메일을 넘어 전반적인 업무 검색이\n가능한 팀 지식 베이스 구축"),
]

rm_x = Inches(0.9)
rm_start_y = Inches(2.3)

for ri, (phase, title, desc) in enumerate(roadmap):
    ry = rm_start_y + ri * Inches(1.0)

    # Phase number pill — completed phase uses green-ish gray
    pill_color = RGBColor(0x16, 0xA3, 0x4A) if phase == "완료" else RED_600
    pill = add_rounded_rect(slide, rm_x, ry, Inches(1.0), Inches(0.3),
                            fill_color=pill_color)
    add_textbox(slide, rm_x, ry + Inches(0.02), Inches(1.0), Inches(0.25),
                phase, font_size=Pt(9), font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, rm_x + Inches(1.2), ry - Inches(0.02), Inches(4), Inches(0.3),
                title, font_size=Pt(15), font_color=WHITE, bold=True)

    # Description
    add_textbox(slide, rm_x + Inches(1.2), ry + Inches(0.3), Inches(4.5), Inches(0.55),
                desc, font_size=Pt(10), font_color=GRAY_400, line_spacing=1.4)

    # Connector line (vertical)
    if ri < len(roadmap) - 1:
        add_rect(slide, rm_x + Inches(0.5), ry + Inches(0.35),
                 Pt(1.5), Inches(0.6), fill_color=RGBColor(0x50, 0x20, 0x20))

# Right side: Thank you
ty_x = Inches(7.5)
ty_y = Inches(2.3)

# Large decorative quote marks
add_textbox(slide, ty_x, ty_y - Inches(0.2), Inches(1), Inches(0.8),
            "\"", font_size=Pt(72), font_color=RED_600, bold=True)

add_textbox(slide, ty_x + Inches(0.3), ty_y + Inches(0.7), Inches(4.5), Inches(1.2),
            "오늘의 이메일이\n내일의 자산이 됩니다.\n팀의 기억을 자동화하세요.",
            font_size=Pt(18), font_color=WHITE, bold=False, italic=True,
            line_spacing=1.6)

# Divider
add_line(slide, ty_x + Inches(0.3), ty_y + Inches(2.2), Inches(3), RED_600, Pt(2))

# Thank you text
add_textbox(slide, ty_x + Inches(0.3), ty_y + Inches(2.5), Inches(4), Inches(0.5),
            "감사합니다",
            font_size=Pt(36), font_color=WHITE, bold=True)

# Contact info
add_textbox(slide, ty_x + Inches(0.3), ty_y + Inches(3.2), Inches(4), Inches(0.6),
            "PM3팀, 박은진\nejpark@artience.com",
            font_size=Pt(11), font_color=GRAY_500, line_spacing=1.6)

# 3 dots at bottom
for i in range(3):
    add_circle(slide, Inches(6.2) + i * Inches(0.3), Inches(6.3),
               Inches(0.1), fill_color=RED_600)


# ─── Save ─────────────────────────────────────────────────────────────
output_path = "TeamWorkHub_v5_clean.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
