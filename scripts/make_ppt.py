"""TeamWorkHub PPT v2 — 프로페셔널 디자인, 화이트+블루 테마."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 16:9 사이즈 ──────────────────────────────────────────────────── #
W = 13.333  # inches
H = 7.5

# ── 컬러 팔레트 ──────────────────────────────────────────────────── #
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_OFFWHITE    = RGBColor(0xF7, 0xF9, 0xFC)
C_SNOW        = RGBColor(0xEF, 0xF4, 0xFB)
C_ICE         = RGBColor(0xDB, 0xEA, 0xFE)
C_SKY         = RGBColor(0xBF, 0xDB, 0xFE)
C_BLUE        = RGBColor(0x3B, 0x82, 0xF6)
C_BLUE_D      = RGBColor(0x25, 0x63, 0xEB)
C_NAVY        = RGBColor(0x1E, 0x40, 0x6E)
C_DEEPNAVY    = RGBColor(0x0C, 0x1D, 0x3A)
C_DARK        = RGBColor(0x1E, 0x29, 0x3B)
C_TEXT        = RGBColor(0x33, 0x3D, 0x4C)
C_SUBTEXT     = RGBColor(0x6B, 0x72, 0x80)
C_LIGHTTEXT   = RGBColor(0x9C, 0xA3, 0xAF)
C_GREEN       = RGBColor(0x10, 0xB9, 0x81)
C_AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
C_RED         = RGBColor(0xEF, 0x44, 0x44)
C_PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
C_TEAL        = RGBColor(0x14, 0xB8, 0xA6)
C_CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)


# ── 유틸 함수 ────────────────────────────────────────────────────── #

def _bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color, border=None, border_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = border_w
    else:
        s.line.fill.background()
    return s


def _rrect(slide, l, t, w, h, color, border=None, border_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = border_w
    else:
        s.line.fill.background()
    # 모서리 반지름 조절 (더 둥글게)
    try:
        s._element.attrib['{http://schemas.microsoft.com/office/drawing/2010/main}' + 'adjLst'] = ''
    except Exception:
        pass
    return s


def _circle(slide, l, t, sz, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _line(slide, l, t, w, color, thickness=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), thickness)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _txt(slide, l, t, w, h, text, sz=18, color=C_TEXT, bold=False,
         align=PP_ALIGN.LEFT, font="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return tb


def _mtxt(slide, l, t, w, h, lines, sz=15, color=C_TEXT,
          spacing=1.3, font="맑은 고딕", bold_first=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = font
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(sz * (spacing - 1) * 2)
    return tb


def _arrow_right(slide, l, t, color=C_SKY):
    """오른쪽 화살표 (chevron)."""
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(l), Inches(t), Inches(0.45), Inches(0.45))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _tag(slide, l, t, text, bg_color=C_ICE, text_color=C_NAVY, sz=11):
    """작은 태그/뱃지."""
    w = len(text) * 0.12 + 0.3
    _rrect(slide, l, t, w, 0.32, bg_color)
    _txt(slide, l, t + 0.02, w, 0.28, text, sz=sz, color=text_color,
         bold=True, align=PP_ALIGN.CENTER)


# ── 슬라이드 빌더 ────────────────────────────────────────────────── #

def s01_cover(prs):
    """표지."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_DEEPNAVY)

    # 배경 장식 — 큰 원들
    _circle(slide, -2.5, -2, 6, C_DARK)
    _circle(slide, 9.5, 3.5, 7, C_DARK)
    _circle(slide, 11, -1.5, 3, RGBColor(0x15, 0x23, 0x3A))

    # 상단 얇은 악센트 라인
    _rect(slide, 0, 0, W, 0.06, C_BLUE)

    # 좌측 세로 악센트 바
    _rect(slide, 1.2, 2.2, 0.06, 2.5, C_BLUE)

    # 메인 타이틀
    _txt(slide, 1.6, 2.0, 8, 1.0, "TeamWorkHub",
         sz=56, color=C_WHITE, bold=True)

    # 서브타이틀
    _txt(slide, 1.6, 3.2, 9, 0.6,
         "이메일 업무 자동화 시스템",
         sz=26, color=C_SKY)

    # 한 줄 설명
    _txt(slide, 1.6, 4.2, 9, 0.5,
         "Gmail에서 메일을 자동 수집하고, AI가 분석하고, Obsidian 업무일지를 만들어줍니다.",
         sz=15, color=C_LIGHTTEXT)

    # 하단 3개 키워드 태그
    keywords = [("Gmail 연동", 1.6), ("Claude AI 분석", 3.3), ("Obsidian 자동 기록", 5.3)]
    for text, x in keywords:
        _rrect(slide, x, 5.2, len(text) * 0.15 + 0.5, 0.38, RGBColor(0x1A, 0x36, 0x5C))
        _txt(slide, x, 5.23, len(text) * 0.15 + 0.5, 0.35, text,
             sz=12, color=C_SKY, align=PP_ALIGN.CENTER, bold=True)

    # 하단 정보
    _txt(slide, 1.6, 6.5, 8, 0.4,
         "Artience Inc.  |  GMPD Data Team  |  2026. 04",
         sz=12, color=C_LIGHTTEXT)


def s02_problem(prs):
    """문제점."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    # 섹션 헤더
    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "PROBLEM", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 8, 0.7,
         "매일 반복되는 이메일 업무, 이대로 괜찮을까요?",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 4개 문제 카드
    problems = [
        ("30분+", "매일 소요", "매일 아침 이메일 확인에\n30분 이상을 소비합니다", C_RED),
        ("Miss", "업무 누락", "중요한 메일을 놓쳐서\n처리가 늦어집니다", C_AMBER),
        ("???", "이력 추적 불가", "회의 때 '그 메일 뭐였지?'\n검색에 또 시간 소요", C_PURPLE),
        ("Copy", "수동 보고", "주간/월간 보고서를\n매번 수동으로 작성합니다", C_SUBTEXT),
    ]

    for i, (icon, title, desc, color) in enumerate(problems):
        x = 1.0 + i * 2.95
        y = 2.2

        # 카드 배경
        _rrect(slide, x, y, 2.65, 4.2, C_WHITE, border=C_CARD_BORDER)

        # 상단 컬러 바
        _rect(slide, x, y, 2.65, 0.06, color)

        # 아이콘 텍스트 (큰 글씨)
        _txt(slide, x, y + 0.5, 2.65, 0.8, icon,
             sz=36, color=color, bold=True, align=PP_ALIGN.CENTER)

        # 제목
        _txt(slide, x, y + 1.5, 2.65, 0.4, title,
             sz=18, color=C_NAVY, bold=True, align=PP_ALIGN.CENTER)

        # 구분선
        _line(slide, x + 0.6, y + 2.1, 1.45, C_CARD_BORDER, Pt(1))

        # 설명
        _txt(slide, x + 0.3, y + 2.4, 2.05, 1.5, desc,
             sz=14, color=C_SUBTEXT, align=PP_ALIGN.CENTER)


def s03_solution(prs):
    """솔루션 개요 — 파이프라인 플로우."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "SOLUTION", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 10, 0.7,
         "TeamWorkHub가 이 모든 걸 자동으로 처리합니다",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 파이프라인 4단계
    steps = [
        ("01", "수집", "Gmail에서\n이메일 자동 수집", C_BLUE),
        ("02", "분석", "Claude AI가\n내용 분석 & 요약", C_BLUE_D),
        ("03", "기록", "Obsidian에\n노트 자동 생성", C_NAVY),
        ("04", "보고", "일간/주간/월간\n보고서 자동 생성", C_GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = 0.8 + i * 3.2
        y = 2.3

        # 번호 원
        _circle(slide, x + 0.7, y, 1.1, color)
        _txt(slide, x + 0.7, y + 0.15, 1.1, 0.5, num,
             sz=28, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 제목
        _txt(slide, x, y + 1.3, 2.5, 0.4, title,
             sz=20, color=C_NAVY, bold=True, align=PP_ALIGN.CENTER)

        # 설명
        _txt(slide, x, y + 1.8, 2.5, 0.8, desc,
             sz=14, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

        # 화살표
        if i < 3:
            _arrow_right(slide, x + 2.5, y + 0.35, C_SKY)

    # 하단 핵심 포인트
    _rrect(slide, 0.8, 5.0, 11.7, 2.0, C_SNOW, border=C_ICE)

    points = [
        "완전 자동화    서버가 정해진 시간에 이메일을 수집하고 정리합니다. 사람이 개입할 필요 없습니다.",
        "중복 처리 방지    한 번 처리된 이메일은 다시 처리하지 않아 파일이 꼬이지 않습니다.",
        "주말도 커버    월요일에는 금요일 저녁부터 일요일 메일까지 빠짐없이 수집합니다.",
    ]
    for i, p in enumerate(points):
        _txt(slide, 1.3, 5.2 + i * 0.55, 10.5, 0.45, p,
             sz=13, color=C_TEXT)
        # 불릿 점
        _circle(slide, 1.1, 5.32 + i * 0.55, 0.12, C_BLUE)


def s04_ai(prs):
    """AI 분석 기능."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "FEATURE 01", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 10, 0.7,
         "AI가 이메일을 읽고 분석합니다",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 5개 분석 항목
    items = [
        ("3줄 요약", "이메일 핵심 내용을\n불릿포인트 3줄로 요약", C_BLUE),
        ("담당자 추출", "본문에서 담당자를 찾고\n닉네임도 풀네임으로 변환", C_TEAL),
        ("긴급도 판단", "긴급 / 보통 / 낮음\n자동 판단", C_RED),
        ("카테고리 분류", "보고 / 승인요청 / 공지\n미팅 / 일반", C_PURPLE),
        ("제목 요약", "긴 메일 제목을\n한국어 핵심 제목으로 압축", C_AMBER),
    ]

    for i, (title, desc, color) in enumerate(items):
        x = 0.5 + i * 2.55
        y = 2.1

        _rrect(slide, x, y, 2.3, 2.6, C_WHITE, border=C_CARD_BORDER)
        # 상단 컬러 악센트
        _rect(slide, x + 0.3, y + 0.25, 1.7, 0.05, color)

        _txt(slide, x, y + 0.5, 2.3, 0.35, title,
             sz=16, color=C_NAVY, bold=True, align=PP_ALIGN.CENTER)
        _txt(slide, x + 0.2, y + 1.1, 1.9, 1.2, desc,
             sz=13, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

    # 하단: RE:RE:RE: 체인 처리
    _rrect(slide, 0.8, 5.1, 11.7, 2.0, C_NAVY)

    _txt(slide, 1.3, 5.25, 5, 0.4,
         "RE: RE: RE:  회신 체인도 똑똑하게 처리",
         sz=18, color=C_WHITE, bold=True)

    _txt(slide, 1.3, 5.75, 10.5, 0.5,
         "이메일 체인에서 최신 답장만 자동 추출하여 AI에게 전달합니다.",
         sz=14, color=C_SKY)

    chain_items = [
        "Outlook / Gmail / 한국어 메일 등 5가지 구분자 패턴 자동 감지",
        "이전 인용 메시지는 맥락 참고만, 최신 답장 중심으로 정확한 요약 생성",
    ]
    for i, item in enumerate(chain_items):
        _circle(slide, 1.3, 6.25 + i * 0.35, 0.08, C_BLUE)
        _txt(slide, 1.6, 6.15 + i * 0.35, 10, 0.35, item,
             sz=12, color=C_LIGHTTEXT)


def s05_daily(prs):
    """Daily Note."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "FEATURE 02", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 10, 0.7,
         "매일 아침, 업무일지가 자동으로 완성됩니다",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 왼쪽: Daily Note 미리보기 (코드 스타일)
    _rrect(slide, 0.8, 2.0, 6.2, 5.0, C_DEEPNAVY)

    # 상단 바 (에디터 느낌)
    _rect(slide, 0.8, 2.0, 6.2, 0.4, C_DARK)
    _circle(slide, 1.0, 2.1, 0.15, C_RED)
    _circle(slide, 1.25, 2.1, 0.15, C_AMBER)
    _circle(slide, 1.5, 2.1, 0.15, C_GREEN)
    _txt(slide, 2.0, 2.05, 4, 0.3, "2026-04-21.md",
         sz=11, color=C_LIGHTTEXT, font="Consolas")

    code_lines = [
        "### Today's work",
        "#### To do list",
        "- [ ] [[미디어성과 측정 표준화]]  #이효수",
        "- [ ] [[어필리에이트 트래킹 권한 요청]]  #이해랑",
        "- [ ] [[확인요청]]  #이기정",
        "- [ ] [[현지화폐 전환 기능 테스트]]  #심민지",
        "- [ ] [[DV360 데이터검증]]  #이기정",
        "- [ ] [[결측치모니터링 로직검수]]  #박은진",
        "      ...(20건 자동 수집)",
        "",
        "#### 정기적인 일",
        "- [ ] 로직점검",
        "",
        "### 미완료  (14일 미처리 항목 자동 표시)",
    ]
    _mtxt(slide, 1.1, 2.6, 5.6, 4.2, code_lines,
          sz=12, color=C_SKY, font="Consolas", spacing=1.25)

    # 오른쪽: 설명 카드들
    features = [
        ("To do list", "전날 저녁 ~ 당일 아침 이메일을\n체크박스 + 위키링크로 자동 생성\n클릭하면 상세 이메일 노트로 이동", C_BLUE),
        ("담당자 태깅", "AI가 이메일에서 담당자 자동 추출\n효수 -> 이효수 (닉네임 자동 변환)\n미지정 시 #미지정 태그 부여", C_TEAL),
        ("정기적인 일", "요일별 정기 업무 자동 표시\n월: RPA / 화: 로직점검 / 수: 수정기\n목: 목정기 / 금: 금정기", C_PURPLE),
        ("미완료 추적", "최근 14일간 체크 안 된 항목을\n자동으로 모아서 표시\nObsidian Dataview 플러그인 연동", C_AMBER),
    ]

    for i, (title, desc, color) in enumerate(features):
        x = 7.5
        y = 2.0 + i * 1.25

        _line(slide, x, y + 0.12, 0.06, color, Pt(35))  # 세로 악센트 바
        _txt(slide, x + 0.2, y, 5, 0.3, title,
             sz=14, color=C_NAVY, bold=True)
        _txt(slide, x + 0.2, y + 0.35, 5, 0.8, desc,
             sz=11, color=C_SUBTEXT)


def s06_reports(prs):
    """주간 / 월간 / 대시보드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "FEATURE 03", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 10, 0.7,
         "보고서도 버튼 하나로 자동 생성",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 3개 보고서 카드
    reports = [
        ("Weekly", "주간 보고서", "2026-W16.md",
         ["한 주간 처리된 이메일 통계",
          "카테고리별 업무 분류",
          "미처리 항목 체크리스트",
          "팀별 업무량 파악 가능"], C_BLUE),
        ("Monthly", "월간 보고서", "2026-04.md",
         ["월간 이메일 처리 현황",
          "발신자 TOP 5 랭킹",
          "담당자별 업무 분포",
          "추이 분석용 데이터"], C_NAVY),
        ("Dashboard", "대시보드", "Dashboard.md",
         ["실시간 업무 현황 집계",
          "담당자별 페이지 자동 생성",
          "Dataview 쿼리 기반",
          "한눈에 팀 현황 파악"], C_GREEN),
    ]

    for i, (tag, title, filename, bullets, color) in enumerate(reports):
        x = 0.8 + i * 4.1
        y = 2.1

        # 카드
        _rrect(slide, x, y, 3.7, 3.5, C_WHITE, border=C_CARD_BORDER)

        # 태그
        _rrect(slide, x + 0.3, y + 0.3, 1.2, 0.32, color)
        _txt(slide, x + 0.3, y + 0.32, 1.2, 0.28, tag,
             sz=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 제목
        _txt(slide, x + 0.3, y + 0.8, 3.1, 0.35, title,
             sz=20, color=C_NAVY, bold=True)

        # 파일명
        _txt(slide, x + 0.3, y + 1.2, 3.1, 0.3, filename,
             sz=11, color=C_SUBTEXT, font="Consolas")

        # 불릿
        for j, b in enumerate(bullets):
            _circle(slide, x + 0.35, y + 1.72 + j * 0.38, 0.07, color)
            _txt(slide, x + 0.55, y + 1.6 + j * 0.38, 2.8, 0.3, b,
                 sz=12, color=C_TEXT)

    # 하단: 본문 가독성 개선
    _rrect(slide, 0.8, 5.9, 11.7, 1.2, C_SNOW, border=C_ICE)
    _txt(slide, 1.3, 6.0, 4, 0.35,
         "이메일 본문 자동 정리",
         sz=16, color=C_NAVY, bold=True)

    cleanup_items = [
        ("제거", "깨진 이미지 태그, 트래킹 픽셀, 이메일 서명, 법적 고지문, 외부메일 경고 배너"),
        ("결과", "Obsidian에서 깔끔한 본문만 표시 — 불필요한 HTML 잔해 없이 가독성 UP"),
    ]
    for i, (label, desc) in enumerate(cleanup_items):
        _rrect(slide, 1.3 + i * 5.5, 6.45, 0.55, 0.28, C_BLUE if i == 1 else C_AMBER)
        _txt(slide, 1.3 + i * 5.5, 6.47, 0.55, 0.24, label,
             sz=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        _txt(slide, 2.0 + i * 5.5, 6.42, 4.8, 0.35, desc,
             sz=11, color=C_SUBTEXT)


def s07_architecture(prs):
    """시스템 구조."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "ARCHITECTURE", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 10, 0.7,
         "시스템 구조",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 아키텍처 블록
    blocks = [
        ("Gmail\nAPI", "이메일 수집\n첨부파일 다운로드", C_BLUE),
        ("Google\nDrive", "첨부파일 업로드\n링크 자동 생성", C_TEAL),
        ("Claude\nAI", "요약 + 담당자\n긴급도 + 카테고리", C_PURPLE),
        ("Obsidian\nVault", "노트 / Daily\nWeekly / Monthly", C_NAVY),
    ]

    for i, (title, desc, color) in enumerate(blocks):
        x = 0.8 + i * 3.2
        y = 2.2

        _rrect(slide, x, y, 2.6, 2.2, color)
        _txt(slide, x, y + 0.3, 2.6, 0.7, title,
             sz=20, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER,
             font="맑은 고딕")
        _txt(slide, x, y + 1.2, 2.6, 0.7, desc,
             sz=12, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

        if i < 3:
            _arrow_right(slide, x + 2.55, y + 0.9, C_SKY)

    # 폴더 구조
    _rrect(slide, 0.8, 4.8, 11.7, 2.3, C_WHITE, border=C_CARD_BORDER)
    _txt(slide, 1.3, 4.95, 5, 0.35,
         "Obsidian Vault 폴더 구조",
         sz=16, color=C_NAVY, bold=True)

    folders = [
        ("TeamWorkHub/", "개별 이메일 노트", "메일 제목.md", C_BLUE),
        ("TeamWorkHub_Daily/", "일간 업무일지", "2026-04-21.md", C_TEAL),
        ("TeamWorkHub_Weekly/", "주간 보고서", "2026-W16.md", C_PURPLE),
        ("TeamWorkHub_Monthly/", "월간 보고서", "2026-04.md", C_NAVY),
        ("TeamWorkHub_Dashboard/", "대시보드 + 담당자", "Dashboard.md", C_GREEN),
    ]

    for i, (name, desc, example, color) in enumerate(folders):
        x = 1.0 + (i % 3) * 3.8
        y = 5.5 + (i // 3) * 1.1

        _line(slide, x, y + 0.08, 0.05, color, Pt(25))
        _txt(slide, x + 0.2, y, 3.2, 0.25, name,
             sz=13, color=C_NAVY, bold=True, font="Consolas")
        _txt(slide, x + 0.2, y + 0.3, 3.2, 0.2, f"{desc}  ({example})",
             sz=11, color=C_SUBTEXT)


def s08_results(prs):
    """기대 효과."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)
    _txt(slide, 1.0, 0.5, 3, 0.35, "IMPACT", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.0, 0.85, 10, 0.7,
         "도입 효과",
         sz=30, color=C_NAVY, bold=True)
    _line(slide, 1.0, 1.55, 2.5, C_BLUE, Pt(3))

    # 4대 지표
    metrics = [
        ("30min+", "매일 절약", "이메일 정리에 쓰던\n시간을 돌려받습니다", C_BLUE),
        ("0건", "업무 누락", "담당자 자동 배정으로\n빠지는 일이 없습니다", C_GREEN),
        ("100%", "기록 보존", "모든 이메일이 자동으로\n기록되고 검색 가능합니다", C_PURPLE),
        ("1-click", "보고서 생성", "일간/주간/월간 보고서를\n자동으로 만들어줍니다", C_TEAL),
    ]

    for i, (num, label, desc, color) in enumerate(metrics):
        x = 0.8 + i * 3.15
        y = 2.1

        _rrect(slide, x, y, 2.8, 3.0, C_WHITE, border=C_CARD_BORDER)

        # 숫자
        _txt(slide, x, y + 0.3, 2.8, 0.7, num,
             sz=38, color=color, bold=True, align=PP_ALIGN.CENTER)

        # 라벨
        _txt(slide, x, y + 1.1, 2.8, 0.35, label,
             sz=16, color=C_NAVY, bold=True, align=PP_ALIGN.CENTER)

        # 구분선
        _line(slide, x + 0.6, y + 1.6, 1.6, C_ICE, Pt(1))

        # 설명
        _txt(slide, x + 0.3, y + 1.8, 2.2, 0.9, desc,
             sz=13, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

    # Before / After
    _rrect(slide, 0.8, 5.5, 5.6, 1.5, RGBColor(0xFE, 0xF2, 0xF2), border=RGBColor(0xFE, 0xCA, 0xCA))
    _txt(slide, 1.2, 5.6, 2, 0.35, "Before", sz=16, color=C_RED, bold=True)
    _mtxt(slide, 1.2, 6.0, 4.8, 0.9, [
        "매일 아침 이메일 30분+ 수동 정리",
        "중요한 메일 놓치고 담당자 불명확",
        "주간/월간 보고서 수동 취합에 또 시간 소요",
    ], sz=12, color=C_TEXT, spacing=1.3)

    _rrect(slide, 6.8, 5.5, 5.7, 1.5, RGBColor(0xEC, 0xFD, 0xF5), border=RGBColor(0xA7, 0xF3, 0xD0))
    _txt(slide, 7.2, 5.6, 2, 0.35, "After", sz=16, color=C_GREEN, bold=True)
    _mtxt(slide, 7.2, 6.0, 5.0, 0.9, [
        "출근하면 업무일지가 이미 완성되어 있음",
        "담당자별 할 일이 명확하게 태깅됨",
        "보고서는 API 한 번 호출로 자동 생성",
    ], sz=12, color=C_TEXT, spacing=1.3)


def s09_next(prs):
    """향후 계획 + Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_DEEPNAVY)

    _rect(slide, 0, 0, W, 0.06, C_BLUE)

    # 장식
    _circle(slide, 10, -1, 4, C_DARK)
    _circle(slide, -1.5, 5, 4, C_DARK)

    _txt(slide, 1.2, 0.6, 3, 0.35, "NEXT STEPS", sz=13, color=C_BLUE, bold=True)
    _txt(slide, 1.2, 0.95, 10, 0.6,
         "향후 계획",
         sz=30, color=C_WHITE, bold=True)
    _line(slide, 1.2, 1.5, 2, C_BLUE, Pt(3))

    plans = [
        ("1", "Cloud Run 배포", "GCP 서버에 올려서 매일 자동 실행  (평일 08:00 KST)", C_BLUE),
        ("2", "Calendar 연동", "Google Calendar 미팅 일정을 Daily Note에 자동 표시", C_TEAL),
        ("3", "실시간 알림", "긴급 메일이 오면 Slack / Teams로 실시간 알림 전송", C_PURPLE),
        ("4", "대시보드 고도화", "팀 전체 업무 현황을 한눈에 볼 수 있는 뷰 구성", C_GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(plans):
        y = 2.0 + i * 0.95

        _circle(slide, 1.2, y + 0.05, 0.45, color)
        _txt(slide, 1.2, y + 0.08, 0.45, 0.35, num,
             sz=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

        _txt(slide, 1.9, y, 4, 0.35, title,
             sz=20, color=C_WHITE, bold=True)
        _txt(slide, 1.9, y + 0.38, 9, 0.35, desc,
             sz=13, color=C_LIGHTTEXT)

        if i < 3:
            _line(slide, 1.42, y + 0.55, 0.01, RGBColor(0x2A, 0x3D, 0x5C), Pt(20))

    # Thank you
    _txt(slide, 1.2, 6.0, 10, 0.7, "Thank you.",
         sz=40, color=C_WHITE, bold=True)
    _txt(slide, 1.2, 6.6, 10, 0.4,
         "TeamWorkHub  |  Powered by Claude AI  |  Artience GMPD Data Team",
         sz=12, color=C_LIGHTTEXT)


# ── 메인 ──────────────────────────────────────────────────────────── #

def main():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    s01_cover(prs)
    s02_problem(prs)
    s03_solution(prs)
    s04_ai(prs)
    s05_daily(prs)
    s06_reports(prs)
    s07_architecture(prs)
    s08_results(prs)
    s09_next(prs)

    out = "TeamWorkHub_소개_v2.pptx"
    prs.save(out)
    print(f"PPT saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
